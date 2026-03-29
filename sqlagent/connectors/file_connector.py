"""File-based connectors: DuckDB, CSV, XLSX, Parquet, JSON.

Files are ingested into DuckDB in-process for SQL querying.
"""

from __future__ import annotations

import os
from datetime import datetime, timezone

import duckdb
import pandas as pd
import structlog

from sqlagent.models import SchemaSnapshot, SchemaTable, SchemaColumn, SampleData, ColumnStats

logger = structlog.get_logger()


class DuckDBConnector:
    """DuckDB in-process connector."""

    def __init__(self, source_id: str, db_path: str = ":memory:"):
        self._source_id = source_id
        self._db_path = db_path
        self._conn = None

    @property
    def dialect(self) -> str:
        return "duckdb"

    @property
    def source_id(self) -> str:
        return self._source_id

    async def connect(self) -> None:
        self._conn = duckdb.connect(self._db_path)

    def _ensure_conn(self):
        if self._conn is None:
            self._conn = duckdb.connect(self._db_path)

    async def execute(self, sql: str, timeout_s: float = 30.0, max_rows: int = 10_000) -> pd.DataFrame:
        self._ensure_conn()
        try:
            result = self._conn.execute(sql)
            df = result.fetchdf()
            return df.head(max_rows)
        except Exception as e:
            from sqlagent.exceptions import SQLExecutionFailed
            raise SQLExecutionFailed(sql=sql, error=str(e))

    async def introspect(self) -> SchemaSnapshot:
        self._ensure_conn()
        tables = []

        result = self._conn.execute("SHOW TABLES")
        table_names = [row[0] for row in result.fetchall()]

        for tname in table_names:
            col_result = self._conn.execute(f"DESCRIBE \"{tname}\"")
            col_rows = col_result.fetchall()
            columns = [
                SchemaColumn(
                    name=r[0],
                    data_type=r[1],
                    nullable=r[2] == "YES" if len(r) > 2 else True,
                    column_position=i,
                )
                for i, r in enumerate(col_rows)
            ]

            try:
                count = self._conn.execute(f"SELECT COUNT(*) FROM \"{tname}\"").fetchone()[0]
            except Exception as exc:
                logger.debug("connector.file.operation_failed", error=str(exc))
                count = 0

            tables.append(SchemaTable(name=tname, columns=columns, row_count_estimate=count))

        return SchemaSnapshot(
            source_id=self._source_id, dialect="duckdb", tables=tables,
        )

    async def sample(self, table: str, n: int = 5) -> SampleData:
        df = await self.execute(f'SELECT * FROM "{table}" LIMIT {n}')
        sample_rows = df.to_dict("records")

        # Compute column stats: distinct values for low-cardinality columns
        col_stats: dict[str, ColumnStats] = {}
        for col in df.columns:
            try:
                dist_df = await self.execute(
                    f'SELECT COUNT(DISTINCT "{col}") AS n FROM "{table}"'
                )
                distinct = int(dist_df.iloc[0, 0]) if not dist_df.empty else 0
                samples: list = []
                if 0 < distinct <= 50:
                    vals_df = await self.execute(
                        f'SELECT DISTINCT "{col}" FROM "{table}" WHERE "{col}" IS NOT NULL LIMIT 30'
                    )
                    samples = [v for v in vals_df.iloc[:, 0].tolist() if v is not None]
                col_stats[col] = ColumnStats(distinct_count=distinct, sample_values=samples)
            except Exception as exc:
                logger.debug("connector.file.operation_failed", error=str(exc))

        return SampleData(table=table, sample_rows=sample_rows, column_stats=col_stats)

    async def health_check(self) -> bool:
        try:
            self._ensure_conn()
            self._conn.execute("SELECT 1")
            return True
        except Exception as exc:
            logger.debug("connector.file.operation_failed", error=str(exc))
            return False

    def register_dataframe(self, name: str, df: pd.DataFrame) -> None:
        """Register a pandas DataFrame as a table in DuckDB."""
        self._ensure_conn()
        self._conn.register(name, df)


class FileConnector(DuckDBConnector):
    """Ingests CSV/XLSX/Parquet/JSON files into DuckDB for SQL querying."""

    def __init__(self, source_id: str, file_path: str, read_config=None):
        """
        Args:
            source_id:   Unique source identifier.
            file_path:   Absolute path to the file.
            read_config: Optional ``ReadConfig`` produced by ``data_profiler.profile_file()``.
                         When provided, its LLM-generated DuckDB arguments are used to load
                         the file cleanly (correct null strings, date formats, cast exprs, etc.).
                         When None, the connector falls back to safe read_csv_auto defaults.
        """
        super().__init__(source_id=source_id, db_path=":memory:")
        self._file_path = file_path
        self._table_name = ""
        self._read_config = read_config  # sqlagent.data_profiler.ReadConfig | None

    async def connect(self) -> None:
        await super().connect()
        await self._ingest_file()

    def _safe_table_name(self) -> str:
        import re
        base = os.path.basename(self._file_path)
        name = re.sub(r'[^a-z0-9_]', '_', os.path.splitext(base)[0].lower()).strip('_')
        return re.sub(r'_+', '_', name)

    async def _ingest_file(self) -> None:
        self._ensure_conn()
        import re
        ext = os.path.splitext(self._file_path)[1].lower()
        self._table_name = self._safe_table_name()

        # ── CSV / TSV / TXT ───────────────────────────────────────────────────
        if ext in (".csv", ".tsv", ".txt"):
            if self._read_config is not None:
                # Use the LLM-generated read_csv_auto() arguments.
                sql = self._read_config.build_read_sql(self._table_name)
                logger.info(
                    "file.ingest_with_profile",
                    file=os.path.basename(self._file_path),
                    null_strings=self._read_config.extra_null_strings,
                    duckdb_args=self._read_config.duckdb_args,
                    cast_exprs=list(self._read_config.cast_exprs.keys()),
                )
                # build_read_sql may return two statements separated by \n
                for stmt in sql.split(";\n"):
                    stmt = stmt.strip()
                    if stmt:
                        self._conn.execute(stmt if stmt.endswith(";") else stmt)
            else:
                # Safe defaults — same behaviour as before, just with better options
                delim_arg = "delim='\\t', " if ext == ".tsv" else ""
                escaped = self._file_path.replace("'", "''")
                self._conn.execute(
                    f"CREATE TABLE \"{self._table_name}\" AS SELECT * FROM "
                    f"read_csv_auto('{escaped}', {delim_arg}"
                    f"normalize_names=true, null_padding=true, ignore_errors=true)"
                )
        elif ext in (".xlsx", ".xls"):
            # Read Excel with pandas (openpyxl for .xlsx, xlrd for .xls)
            engine = "openpyxl" if ext == ".xlsx" else "xlrd"
            try:
                # Try reading all sheets
                xls = pd.ExcelFile(self._file_path, engine=engine)
                def _sanitize_name(name):
                    """Make a string safe for SQL table names."""
                    import re
                    return re.sub(r'[^a-z0-9_]', '_', name.lower()).strip('_')

                if len(xls.sheet_names) == 1:
                    df = pd.read_excel(xls, sheet_name=0)
                    safe = _sanitize_name(self._table_name)
                    self._table_name = safe
                    self._conn.register(safe, df)
                    self._conn.execute(
                        f'CREATE TABLE "{safe}" AS SELECT * FROM "{safe}"'
                    )
                else:
                    # Multiple sheets → one table per sheet
                    for sheet in xls.sheet_names:
                        safe = _sanitize_name(self._table_name + "_" + sheet)
                        df = pd.read_excel(xls, sheet_name=sheet)
                        if df.empty:
                            continue
                        self._conn.register(safe, df)
                        self._conn.execute(
                            f'CREATE TABLE "{safe}" AS SELECT * FROM "{safe}"'
                        )
            except Exception as e:
                raise ValueError(f"Failed to read Excel file: {e}")
        elif ext == ".parquet":
            self._conn.execute(
                f"CREATE TABLE \"{self._table_name}\" AS SELECT * FROM read_parquet('{self._file_path}')"
            )
        elif ext == ".json":
            escaped = self._file_path.replace("'", "''")
            self._conn.execute(
                f"CREATE TABLE \"{self._table_name}\" AS SELECT * FROM "
                f"read_json_auto('{escaped}', ignore_errors=true)"
            )
        elif ext in (".pptx", ".ppt"):
            # Extract text/tables from PowerPoint slides into a table
            try:
                from pptx import Presentation
                prs = Presentation(self._file_path)
                rows = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    for shape in slide.shapes:
                        if shape.has_table:
                            tbl = shape.table
                            headers = [cell.text for cell in tbl.rows[0].cells]
                            for row in tbl.rows[1:]:
                                row_data = {"slide": slide_num}
                                for j, cell in enumerate(row.cells):
                                    col_name = headers[j] if j < len(headers) else f"col_{j}"
                                    row_data[col_name] = cell.text
                                rows.append(row_data)
                        elif shape.has_text_frame:
                            text = shape.text_frame.text.strip()
                            if text:
                                rows.append({"slide": slide_num, "content": text, "shape_type": "text"})
                if rows:
                    df = pd.DataFrame(rows)
                    self._conn.register(self._table_name, df)
                    self._conn.execute(
                        f"CREATE TABLE \"{self._table_name}\" AS SELECT * FROM {self._table_name}"
                    )
                else:
                    raise ValueError("No extractable content found in PowerPoint file")
            except ImportError:
                raise ValueError("python-pptx required for PowerPoint files: pip install python-pptx")
        else:
            raise ValueError(f"Unsupported file format: {ext}")

        logger.info("file.ingested", file=self._file_path, table=self._table_name)
